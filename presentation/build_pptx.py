#!/usr/bin/env python3
"""SIGMA 소개 덱 생성기 — 미니멀 16:9, 실제 화면 중심, 15장(상세판)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
SLIDES = os.path.join(HERE, "slides")

INK = RGBColor(0x1C, 0x1B, 0x19)
GRAY = RGBColor(0x6F, 0x6D, 0x68)
FAINT = RGBColor(0xA8, 0xA5, 0x9D)
BG = RGBColor(0xFA, 0xF9, 0xF7)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE2, 0xDF, 0xD9)
FONT = "Apple SD Gothic Neo"

EMU_IN = 914400
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Emu(int(SW * EMU_IN))
prs.slide_height = Emu(int(SH * EMU_IN))
BLANK = prs.slide_layouts[6]


def _set_font(run, name):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def text(s, txt, left, top, width, height, size, *, bold=False, color=INK,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None, font=FONT):
    box = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(txt.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        _set_font(r, font)
    return box


def picture(s, path, max_w, max_h, top, left=None, border=True):
    w, h = Image.open(path).size
    r = min(max_w / w, max_h / h)
    wi, hi = w * r, h * r
    if left is None:
        left = (SW - wi) / 2
    pic = s.shapes.add_picture(path, Inches(left), Inches(top), Inches(wi), Inches(hi))
    if border:
        pic.line.color.rgb = BORDER
        pic.line.width = Pt(0.75)
    return pic


def page_no(s, n):
    text(s, f"{n}", SW - 1.0, SH - 0.5, 0.6, 0.3, 10, color=FAINT, align=PP_ALIGN.RIGHT)


def card(s, x, y, w, h, title, desc, title_sz=18, desc_sz=12.5):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = CARD
    box.line.color.rgb = BORDER; box.line.width = Pt(1)
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title; r.font.size = Pt(title_sz); r.font.bold = True; r.font.color.rgb = INK; _set_font(r, FONT)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.line_spacing = 1.15
    r2 = p2.add_run(); r2.text = desc; r2.font.size = Pt(desc_sz); r2.font.color.rgb = GRAY; _set_font(r2, FONT)


def feature(n, title, bullets, fname):
    s = slide()
    text(s, title, 0.75, 0.5, SW - 1.5, 0.7, 26, bold=True, color=INK)
    text(s, "\n".join("•  " + b for b in bullets), 0.78, 1.28, SW - 1.5, 1.2, 14.5,
         color=INK, spacing=1.28)
    path = os.path.join(SLIDES, fname)
    if os.path.exists(path):
        picture(s, path, SW - 1.6, 4.35, 2.62)
    page_no(s, n)


# ── 1. 타이틀 ───────────────────────────────────────────────────────────
s = slide()
text(s, "SIGMA", 0, 2.45, SW, 1.4, 60, bold=True, color=INK, align=PP_ALIGN.CENTER)
text(s, "지금 할 일을, 알아서 위로.", 0, 3.85, SW, 0.6, 22, color=GRAY, align=PP_ALIGN.CENTER)
text(s, "적응형 할 일 관리", 0, 4.5, SW, 0.4, 14, color=FAINT, align=PP_ALIGN.CENTER)

# ── 2. 문제 ─────────────────────────────────────────────────────────────
s = slide()
text(s, "할 일이 쌓이면,\n가장 어려운 건 “지금 뭘 하지?”", 0, 2.2, SW, 1.8, 32, bold=True,
     color=INK, align=PP_ALIGN.CENTER, spacing=1.25)
text(s, "마감 · 중요도 · 미룬 일 … 매번 머릿속으로 다시 계산해야 한다.", 0, 4.25, SW, 0.6, 16,
     color=GRAY, align=PP_ALIGN.CENTER)
page_no(s, 2)

# ── 3. 핵심 한눈에 ──────────────────────────────────────────────────────
s = slide()
text(s, "SIGMA는, 세 가지로", 0.75, 0.7, SW - 1.5, 0.7, 27, bold=True, color=INK)
items = [
    ("01  적응형 우선순위", "지금 할 수 있는 일을 위로 — ‘최우선 과제’ 한 가지를 제시"),
    ("02  키보드 우선", "⌘K 커맨드 팔레트 · N 빠른 추가로 검색·이동·생성을 한 번에"),
    ("03  놓침 방지", "여러 번 미룬 ‘좀비’·마감 임박을 알림으로, 보관함으로 안전하게"),
]
y = 2.1
for t, d in items:
    text(s, t, 1.1, y, 4.6, 0.5, 20, bold=True, color=INK)
    text(s, d, 5.9, y + 0.03, 6.3, 0.8, 15, color=GRAY, spacing=1.1)
    y += 1.45
page_no(s, 3)

# ── 4~12. 기능 슬라이드 ─────────────────────────────────────────────────
feature(4, "적응형 우선순위", [
    "가용 시간·에너지를 입력하면 지금 할 수 있는 일이 위로",
    "‘최우선 과제’ 한 가지를 카드로 제시 + 집중 타이머",
    "맞지 않는 작업은 흐리게(시간·에너지 초과)",
], "01-today.png")

feature(5, "한눈에 보는 리스트 · 정렬 · 보기", [
    "오늘 · 중요 · 기록 · 보관함 — 평평한 줄로 빠르게 스캔",
    "정렬: 적응형 · 마감 임박순 · 중요도순  /  보기: 카드 ⇄ 줄",
    "카테고리 의미 색 태그 + D-day 우측 정렬",
], "05-important.png")

feature(6, "빠른 추가", [
    "N 키 또는 ‘+ 새 작업’으로 즉시 추가",
    "제목 + 에너지 · 카테고리 · 중요도 · 마감을 한 번에",
    "추가 직후 우선순위에 자동 반영",
], "08-composer.png")

feature(7, "⌘K 커맨드 팔레트", [
    "⌘K / Ctrl+K — 어디서나 호출",
    "모든 작업(오늘·보관함·기록) 통합 검색 + 페이지 이동",
    "‘…’ 새 작업으로 추가까지 한 곳에서, 키보드만으로",
], "02-palette.png")

feature(8, "놓치지 않게 — 좀비 작업", [
    "5번 이상 미룬 작업은 빨갛게 ‘좀비’로 강조",
    "‘지금은 타이밍이 아닐 수 있어요’ — 보관 / 계속 유지 제안",
    "보관하면 보관함으로, 언제든 복구",
], "09-zombie.png")

feature(9, "보관함 — 복구 · 영구삭제", [
    "미뤄둔 작업을 한곳에 모아 시야에서 정리",
    "‘복구’로 다시 목록에, ‘삭제’로 영구 제거(5초 실행취소)",
    "보관 시 토스트로 어디로 갔는지 즉시 안내",
], "03-archive.png")

feature(10, "알림", [
    "여러 번 미룬 작업 · 마감 임박을 벨 알림으로",
    "놓치기 쉬운 일을 먼저 띄움",
    "설정에서 알림 표시 on/off",
], "04-notif.png")

# ── 11. 미니 위젯 (세로 이미지 → 우측 배치) ────────────────────────────
s = slide()
text(s, "어디서나, 작게", 0.75, 0.5, SW - 1.5, 0.7, 26, bold=True, color=INK)
text(s, "메인 창을 내려도 떠 있는 작은 위젯으로 오늘 할 일을 확인.", 0.78, 1.28, 7.0, 0.5, 14.5, color=GRAY)
wpath = os.path.join(SLIDES, "07-widget.png")
if os.path.exists(wpath):
    picture(s, wpath, 3.0, 4.5, 2.2, left=8.9)
text(s, "•  Chrome · Edge\n     화면 위에 항상 떠 있는 미니 창(PiP)\n\n•  그 외 브라우저\n     화면 오른쪽에 작은 팝업 창\n\n•  위젯에서 체크하면 본 화면도 즉시 동기화",
     0.95, 2.55, 7.4, 3.6, 16, color=INK, spacing=1.3)
page_no(s, 11)

# ── 12. 설정 · 도움말 ───────────────────────────────────────────────────
feature(12, "설정 · 도움말", [
    "테마: 시스템 · 라이트 · 다크 (선택 저장)",
    "알림 표시 토글",
    "도움말: 단축키·기능 안내(⌘K · N · Enter · Esc)",
], "10-settings.png")

# ── 13. 속 엔진 ① 점수 모델 ────────────────────────────────────────────
s = slide()
text(s, "속 엔진 ① — 점수 모델", 0.75, 0.5, SW - 1.5, 0.7, 26, bold=True, color=INK)
text(s, "중요한 일을 자동으로 위로 올리는 우선순위 점수.", 0.78, 1.28, SW - 1.5, 0.5, 14.5, color=GRAY)
text(s, "점수  =  중요도  +  긴급도  −  지연", 0.75, 2.15, SW - 1.5, 0.7, 26, bold=True,
     color=INK, align=PP_ALIGN.CENTER)
cw, ch, gap = 3.3, 1.5, 0.55
x0 = (SW - (cw * 3 + gap * 2)) / 2
for i, (t, d) in enumerate([("중요도", "중요할수록 ↑"), ("긴급도", "마감 임박할수록 ↑"), ("지연", "미룰수록 ↓")]):
    card(s, x0 + i * (cw + gap), 3.25, cw, ch, t, d)
text(s, "긴급도는 마감이 가까울수록 1에 수렴(쌍곡 감쇠) · 마감이 없으면 미룬 횟수로 ‘방치도’가 오른다.",
     0.75, 5.2, SW - 1.5, 0.7, 13.5, color=GRAY, align=PP_ALIGN.CENTER, spacing=1.2)
page_no(s, 13)

# ── 14. 속 엔진 ② 학습 · 탐색 ──────────────────────────────────────────
s = slide()
text(s, "속 엔진 ② — 스스로 적응", 0.75, 0.5, SW - 1.5, 0.7, 26, bold=True, color=INK)
text(s, "행동을 학습해 가중치를 조정하고, 가끔은 새로운 일도 권한다.", 0.78, 1.28, SW - 1.5, 0.5, 14.5, color=GRAY)
cw2, ch2, gap2 = 5.4, 2.0, 0.6
x0 = (SW - (cw2 * 2 + gap2)) / 2
card(s, x0, 2.5, cw2, ch2, "행동 학습",
     "중요한 일을 자꾸 미루면(스누즈↑) 중요도 가중치를 자동 상향.\n매일 자정, 최근 24시간 행동 로그를 학습.", desc_sz=13)
card(s, x0 + cw2 + gap2, 2.5, cw2, ch2, "탐색 · ε-greedy 5%",
     "가끔(5%) 최근 가장 덜 한 카테고리의 일을 슬쩍 추천.\n한쪽으로 치우치는 ‘편식’을 막는다.", desc_sz=13)
text(s, "규칙 기반 · 설명가능 — 왜 이 점수인지(중요도/긴급도/지연 기여분)를 그대로 드러낸다.",
     0.75, 5.0, SW - 1.5, 0.7, 13.5, color=GRAY, align=PP_ALIGN.CENTER, spacing=1.2)
page_no(s, 14)

# ── 15. 마무리 ──────────────────────────────────────────────────────────
s = slide()
text(s, "정리가 아니라, ‘지금 할 일’로.", 0, 2.7, SW, 1.0, 30, bold=True, color=INK, align=PP_ALIGN.CENTER)
text(s, "SIGMA", 0, 3.9, SW, 0.7, 26, bold=True, color=FAINT, align=PP_ALIGN.CENTER)

out = os.path.join(HERE, "SIGMA_소개.pptx")
prs.save(out)
print("saved:", out, "/ slides:", len(prs.slides._sldIdLst))
